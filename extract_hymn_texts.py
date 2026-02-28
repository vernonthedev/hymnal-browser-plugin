import os
from pptx import Presentation

SOURCE_DIR = r"D:\livestreamEKC\sdaFamily\Converted_PowerPoints"
OUTPUT_DIR = r"D:\livestreamEKC\hymn_texts"

os.makedirs(OUTPUT_DIR, exist_ok=True)

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = []

    for slide in prs.slides:
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)

        if slide_lines:
            slides_text.append(" ".join(slide_lines))

    return slides_text


def main():
    files = sorted(f for f in os.listdir(SOURCE_DIR) if f.lower().endswith(".pptx"))

    print("\n====== HYMN TEXT EXTRACTION ======\n")
    print(f"Source: {SOURCE_DIR}")
    print(f"Output: {OUTPUT_DIR}")
    print(f"Files found: {len(files)}\n")

    success = 0

    for file in files:
        hymn_number = os.path.splitext(file)[0]  # e.g. 25.pptx → 25
        pptx_path = os.path.join(SOURCE_DIR, file)

        print(f"[+] Extracting hymn {hymn_number}")

        try:
            slides_text = extract_text_from_pptx(pptx_path)

            if not slides_text:
                print(f"[!] No text found in {file}")
                continue

            out_path = os.path.join(OUTPUT_DIR, f"{hymn_number}.txt")

            with open(out_path, "w", encoding="utf-8") as f:
                for i, line in enumerate(slides_text, 1):
                    f.write(f"[Slide {i}]\n")
                    f.write(line + "\n\n")

            success += 1

        except Exception as e:
            print(f"[ERROR] {file}: {e}")

    print("\n=================================")
    print(f"Completed: {success}/{len(files)} hymns extracted")
    print("=================================\n")


if __name__ == "__main__":
    main()